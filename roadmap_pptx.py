"""
Populate the RoadMap template with extracted charts and values.
Uses placeholder names only: text tokens (e.g. {{LIQUID_ASSETS_PRE}}) and shape names (Selection Pane).
No slide numbers — replacements run across the whole deck so content can be moved or duplicated.
"""

import logging
from pathlib import Path
from typing import Optional, Tuple, Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


def _fmt_liquid_millions(value: Optional[int]) -> str:
    """Format liquid total (in £) as c.£Xm. One decimal if needed; no trailing .0 if whole."""
    if value is None:
        return "—"
    millions = value / 1_000_000
    if millions == int(millions):
        return f"c.£{int(millions)}m"
    return f"c.£{millions:.1f}m"


def _fmt_gbp(value: Optional[int]) -> str:
    """Format int as £ with commas."""
    if value is None:
        return "—"
    return f"£{value:,}"


def _get_slide_text(slide) -> str:
    """Collect all text from a slide (for placeholder check)."""
    def collect(shapes):
        parts = []
        for shape in shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            parts.append(run.text)
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                parts.extend(collect(shape.shapes))
        return parts
    return "".join(collect(slide.shapes))


def _find_shape_by_name(slide, name: str) -> Tuple[Optional[Any], int, int]:
    """
    Find a shape by name on the slide, including inside groups.
    Returns (shape, abs_left, abs_top) in slide coordinates, or (None, 0, 0).
    """
    def search(shapes, offset_left: int, offset_top: int):
        for shape in shapes:
            if shape.name == name:
                return shape, offset_left + shape.left, offset_top + shape.top
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                found, abs_left, abs_top = search(shape.shapes, offset_left + shape.left, offset_top + shape.top)
                if found is not None:
                    return found, abs_left, abs_top
        return None, 0, 0
    return search(slide.shapes, 0, 0)


def _replace_shape_with_image(slide, shape_name: str, image_path: Path) -> bool:
    """
    Replace the first shape with the given name with the image. Uses scale-to-fit so the
    image keeps its aspect ratio and is not stretched; optionally avoids upscaling so charts stay sharp.
    Returns True if replaced.
    """
    shape, abs_left, abs_top = _find_shape_by_name(slide, shape_name)
    if shape is None or not image_path.exists():
        return False
    box_w = shape.width
    box_h = shape.height
    sp = shape._element.getparent()
    sp.remove(shape._element)

    # Get image size in pixels to preserve aspect ratio and avoid stretching
    try:
        with Image.open(image_path) as img:
            img_w_px, img_h_px = img.size
    except Exception:
        img_w_px = img_h_px = None

    if img_w_px and img_h_px and img_w_px > 0 and img_h_px > 0:
        # Scale to fit inside placeholder (EMU per pixel). python-pptx uses EMU; 1 inch = 914400 EMU.
        scale_emu_per_px = min(box_w / img_w_px, box_h / img_h_px)
        # Cap so we never upscale beyond ~96 DPI equivalent (keeps charts sharp)
        emu_per_px_96dpi = 914400 / 96
        scale_emu_per_px = min(scale_emu_per_px, emu_per_px_96dpi)
        new_width = int(img_w_px * scale_emu_per_px)
        new_height = int(img_h_px * scale_emu_per_px)
        # Center image in the placeholder box
        left = abs_left + (box_w - new_width) // 2
        top = abs_top + (box_h - new_height) // 2
        slide.shapes.add_picture(str(image_path), left, top, new_width, new_height)
    else:
        # Fallback: use placeholder dimensions if we couldn't read image size
        slide.shapes.add_picture(str(image_path), abs_left, abs_top, box_w, box_h)
    return True


def _replace_text_tokens_in_shapes(shapes, tokens: dict[str, str]) -> None:
    """Replace tokens in all text in shapes (and recurse into groups)."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            _replace_text_tokens_in_shapes(shape.shapes, tokens)
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in para.runs)
            if not full_text:
                continue
            new_text = full_text
            for token, value in tokens.items():
                new_text = new_text.replace(token, value)
            if new_text == full_text:
                continue
            if para.runs:
                para.runs[0].text = new_text
                for i in range(1, len(para.runs)):
                    para.runs[i].text = ""


def _replace_text_tokens(slide, tokens: dict[str, str]) -> None:
    """Replace placeholder tokens with values in all text on the slide (including inside groups)."""
    _replace_text_tokens_in_shapes(slide.shapes, tokens)


# Shape name -> chart key in all_charts. Every occurrence of each shape name in the deck is replaced.
IMAGE_PLACEHOLDERS = [
    ("[TIMELINE_IMAGE]", "pre_timeline"),
    ("PRE_CASHFLOW_IMAGE", "pre_cashflow"),
    ("PRE_LIQUID_IMAGE", "pre_liquid_assets"),
    ("PRE_CASHFLOW_COMPARISON", "pre_cashflow"),
    ("POST_CASHFLOW_COMPARISON", "post_cashflow"),
    ("PRE_LIQUID_COMPARISON", "pre_liquid_assets"),
    ("POST_LIQUID_COMPARISON", "post_liquid_assets"),
    ("COMP_CHART_1", "slide19_comparison_chart_1"),
    ("COMP_CHART_2", "slide19_comparison_chart_2"),
    ("COMP_CHART_3", "slide19_comparison_chart_3"),
    ("COMP_CHART_4", "slide19_comparison_chart_4"),
    ("POST_ESTATE_IMAGE", "slide24_estate_analysis"),
]


def populate_roadmap_pptx(
    template_path: Path,
    output_path: Path,
    charts_dir: Path,
    all_charts: dict[str, str],
    retirement_annual_diff: int,
    retirement_monthly_diff: int,
    liquid_pre: Optional[int] = None,
    liquid_post: Optional[int] = None,
    shortfall_years: Optional[int] = None,
    total_retirement_years: Optional[int] = None,
    lump_sum_required: Optional[int] = None,
    retirement_year: Optional[int] = None,
    annual_savings_required: Optional[int] = None,
    post_not_funded_years: Optional[int] = None,
    post_funded_years: Optional[int] = None,
    post_retirement_spending: Optional[int] = None,
    template_type: str = "Generic",
    client_name: Optional[str] = None,
    report_month: Optional[str] = None,
    report_year: Optional[str] = None,
) -> None:
    """
    Replace all placeholders in the deck by name. No slide numbers — every occurrence
    of each text token or shape name is replaced, so content can appear anywhere or be duplicated.

    charts_dir is the charts folder (output_dir / "charts").
    all_charts maps key -> relative path e.g. "charts/pre_timeline_page4.png".
    template_type is kept for API compatibility but no longer affects slide indices.
    """
    prs = Presentation(str(template_path))
    base_dir = charts_dir.parent

    def chart_path(key: str) -> Optional[Path]:
        if key not in all_charts:
            return None
        return base_dir / all_charts[key]

    # ---- All text tokens — replace on every slide (so e.g. Slide 6 figures can appear elsewhere too) ----
    pre_funded_years: Optional[int] = None
    if total_retirement_years is not None and shortfall_years is not None:
        pre_funded_years = total_retirement_years - shortfall_years

    all_tokens = {
        "{{CLIENT_NAME}}": (client_name or "—").strip(),
        "{{REPORT_MONTH}}": report_month or "—",
        "{{REPORT_YEAR}}": report_year or "—",
        "{{RETIREMENT_MONTHLY_DIFF}}": f"£{retirement_monthly_diff:,}",
        "{{RETIREMENT_ANNUAL_DIFF}}": f"£{retirement_annual_diff:,}",
        "{{LIQUID_ASSETS_PRE}}": _fmt_liquid_millions(liquid_pre),
        "{{LIQUID_ASSETS_POST}}": _fmt_liquid_millions(liquid_post),
        "{{SHORTFALL_YEARS}}": str(shortfall_years) if shortfall_years is not None else "—",
        "{{TOTAL_RETIREMENT_YEARS}}": str(total_retirement_years) if total_retirement_years is not None else "—",
        "{{PRE_FUNDED_YEARS}}": str(pre_funded_years) if pre_funded_years is not None else "—",
        "{{LUMP_SUM_REQUIRED}}": _fmt_gbp(lump_sum_required),
        "{{RETIREMENT_YEAR}}": str(retirement_year) if retirement_year is not None else "—",
        "{{ANNUAL_SAVINGS_REQUIRED}}": _fmt_gbp(annual_savings_required),
        "{{POST_NOT_FUNDED_YEARS}}": str(post_not_funded_years) if post_not_funded_years is not None else "—",
        "{{POST_FUNDED_YEARS}}": str(post_funded_years) if post_funded_years is not None else "—",
        "{{POST_RETIREMENT_SPENDING}}": _fmt_gbp(post_retirement_spending),
    }

    for slide in prs.slides:
        _replace_text_tokens(slide, all_tokens)

    # Warn if any placeholder token is still present (e.g. typo in template)
    full_text = "".join(_get_slide_text(s) for s in prs.slides)
    still_present = [t for t in all_tokens if t in full_text]
    if still_present:
        logger.warning(
            "Placeholder token(s) still present in deck (not found or not replaced): %s",
            ", ".join(still_present),
        )

    # ---- Image placeholders — replace every occurrence of each shape name across all slides ----
    for shape_name, chart_key in IMAGE_PLACEHOLDERS:
        path = chart_path(chart_key)
        if not path or not path.exists():
            continue
        for slide in prs.slides:
            while _replace_shape_with_image(slide, shape_name, path):
                pass

    prs.save(str(output_path))
